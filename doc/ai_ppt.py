import collections.abc
from pptx import Presentation
from pptx.util import Inches, Pt
from pptx.enum.text import PP_ALIGN

# 프레젠테이션 객체 생성
prs = Presentation()

# 공통 함수
def add_content_slide(prs, title, content_lines, note=None):
    layout = prs.slide_layouts[1]
    slide = prs.slides.add_slide(layout)
    slide.shapes.title.text = title
    
    tf = slide.placeholders[1].text_frame
    tf.clear()

    for i, line in enumerate(content_lines):
        if i == 0:
            p = tf.paragraphs[0]
        else:
            p = tf.add_paragraph()
        
        clean_line = line.strip()
        if clean_line.startswith("-"):
            p.text = clean_line[1:].strip()
            p.level = 1
        elif clean_line.startswith("1.") or clean_line.startswith("2.") or clean_line.startswith("3.") or clean_line.startswith("Step"):
             p.text = clean_line
             p.level = 0
             p.font.bold = True
        elif clean_line == "":
            p.text = ""
        else:
            p.text = clean_line
            p.level = 0

    if note:
        slide.notes_slide.notes_text_frame.text = note
    return slide

# --- 1~10번 슬라이드는 기존 논리 유지 ---

# 1. Intro
slide = prs.slides.add_slide(prs.slide_layouts[0])
slide.shapes.title.text = "개발자의 진화: Coder to Architect"
slide.placeholders[1].text = "Know-How → Know-Where → AI Design\nAI Agent 시대, 우리는 무엇을 준비해야 하는가?"
slide.notes_slide.notes_text_frame.text = "안녕하세요. 오늘 저는 개발자라는 직업이 어떻게 변화해왔고, Cline과 같은 AI 에이전트 시대에 우리는 무엇을 준비해야 하는지에 대해 이야기하려 합니다."

# 2. 화두
add_content_slide(prs, "화두 던지기: 실력의 기준이 바뀌다", 
                  ["\"코딩 잘한다\"의 기준 변화", "- 10년 전: 암기 (알고리즘)", "- 5년 전: 정보 (라이브러리)", "- 지금: ???", "", "AI가 스스로 코드를 고치는 시대, 과연 '실력'은 무엇일까요?"],
                  "여러분은 '개발을 잘한다'는 것이 무엇이라고 생각하시나요? 과거에는 암기, 최근에는 검색이었습니다. AI가 스스로 디버깅까지 하는 지금, 과연 무엇이 실력일까요?")

# 3. Roadmap
add_content_slide(prs, "변화의 흐름 (The 3 Eras)", 
                  ["1. Era 1: Know-How (노하우)", "2. Era 2: Know-Where (노우웨어)", "3. Era 3: AI Design (AI 디자인)"],
                  "저는 개발의 역사를 크게 세 가지 시대로 구분합니다. 노하우, 노우웨어, 그리고 AI 디자인의 시대입니다.")

# 4. Era 1
add_content_slide(prs, "Era 1: Know-How (노하우)", 
                  ["The Craftsman (장인)", "- 지식의 내재화 (Memory)", "- 최적화, 문법 숙지", "", "\"내 머릿속에 모든 코드가 있다.\""],
                  "첫 번째는 '노하우'의 시대입니다. 문법을 외우고 최적화 기법을 손끝에 익히는 '장인 정신'이 핵심이었죠.")

# 5. Era 2
add_content_slide(prs, "Era 2: Know-Where (노우웨어)", 
                  ["The Researcher (검색가/조립가)", "- 지식의 연결 (Link)", "- Google, StackOverflow, GitHub", "", "\"내가 몰라도 인터넷 어딘가에 답이 있다.\""],
                  "두 번째는 '노우웨어'의 시대입니다. 정보를 잘 찾고 조립하는 능력이 중요해졌죠.")

# 6. Era 3
add_content_slide(prs, "Era 3: AI Design (AI 디자인)", 
                  ["The Architect (설계자/지휘자)", "- 지식의 생성 및 지휘 (Direct & Design)", "- ChatGPT, Copilot, Cline (AI Agents)", "", "\"구현과 수정은 AI가 한다. 나는 무엇을 만들지 결정한다.\""],
                  "이제 'AI 디자인'의 시대입니다. AI Agent는 구현뿐만 아니라 수정까지 수행합니다. 이제 중요한 건 '설계'와 '결정'입니다.")

# 7. Deep Dive
add_content_slide(prs, "왜 'Design' 인가?", 
                  ["Coding is Free, Thinking is Expensive", "", "개발자의 역할 변화", "- Writer (작가) → Editor-in-Chief (편집장)", "", "핵심 가치", "- 전체 구조 설계 (Architecture)", "- AI 결과물의 감리 (Auditing)"],
                  "AI는 훌륭한 '작가'이자 '수정가'입니다. 하지만 방향을 잡아줄 '편집장'이 필요합니다. 그것이 바로 디자인입니다.")

# 8. Table Slide
slide = prs.slides.add_slide(prs.slide_layouts[5])
slide.shapes.title.text = "시대별 비교 (Summary)"
rows, cols = 5, 4
left, top, width, height = Inches(0.5), Inches(1.5), Inches(9.0), Inches(0.8)
table = slide.shapes.add_table(rows, cols, left, top, width, height).table
headers = ["구분", "Know-How (과거)", "Know-Where (현재)", "AI Design (미래)"]
for i, h in enumerate(headers): table.cell(0, i).text = h
data = [["핵심 역량", "문법 암기", "검색/조립", "설계/감리/지휘"], ["역할", "장인 (Writer)", "조립가 (Assembler)", "감독 (Director)"], ["작업 방식", "Typing", "Searching", "Designing & Auditing"], ["생산성", "숙련도", "정보력", "AI 협업력"]]
for r, row in enumerate(data):
    for c, val in enumerate(row):
        table.cell(r+1, c).text = val
slide.notes_slide.notes_text_frame.text = "정리하자면 이렇습니다. 과거의 우리가 장인이었고, 현재가 사서였다면, 미래의 우리는 AI 에이전트 팀을 지휘하는 감독이 되어야 합니다."

# 9. Skills
add_content_slide(prs, "새로운 핵심 역량", 
                  ["1. 질문 능력 (Prompt Engineering)", "2. 안목 (Insight)", "3. 시스템 사고 (System Thinking)"],
                  "AI를 지휘할 질문 능력, 결과물을 판단할 안목, 숲을 보는 시스템 사고가 필수적입니다.")

# 10. Action Plan 1
add_content_slide(prs, "Action Plan 1: 설계를 위한 '시각화'와 '분해'", 
                  ["핵심: 코딩하기 전에 그리고, AI가 이해하도록 쪼개라", "", "1. Visual Thinking (시각화)", "- 텍스트보다 다이어그램(Flowchart)으로 먼저 구조 정의", "- AI에게 '코드 짜줘'가 아닌 '이 설계도대로 구현해줘'라고 지시", "", "2. Decomposition (문제 분해)", "- 복잡한 시스템을 AI Agent가 처리 가능한 '단위 모듈'로 분해", "- 높은 응집도와 낮은 결합도 설계가 AI 성능을 좌우함"],
                  "무턱대고 AI에게 시키기 전에, 먼저 설계도를 그려야 합니다. 그리고 Agent가 잘 소화할 수 있도록 문제를 작게 쪼개주는 것이 핵심입니다.")

# --- 수정된 파트: Agent 트렌드 반영 (Slide 11, 12) ---

# 11. Action Plan 2 (Updated)
add_content_slide(prs, "Action Plan 2: 검증을 넘어 '감리(Auditing)'로", 
                  ["핵심: 코드 리뷰(Review)가 아닌 디시전 리뷰(Decision Review)", "", "1. From Syntax to Context", "- 단순 에러/버그는 AI Agent가 스스로 수정(Self-Correction)함", "- 인간은 '비즈니스 로직'과 '방향성'이 맞는지 확인해야 함", "", "2. Human-in-the-loop (중간 승인)", "- AI가 잘못된 가정으로 무한 루프에 빠지지 않도록 방향타 조정", "- 보안 및 유지보수성(Technical Debt)에 대한 '최종 승인' 권한 행사"],
                  "요즘 Cline 같은 에이전트는 스스로 에러를 고칩니다. 이제 우리는 오타를 찾는 게 아니라, AI가 '올바른 문제'를 풀고 있는지, '위험한 방식'을 쓰지 않는지 '감리(Audit)'해야 합니다.")

# 12. Workflow (Updated)
add_content_slide(prs, "The New Workflow: Manager of Agents", 
                  ["Process: Design → Direct → Audit", "", "Step 1: Design (설계)", "- 요구사항 정의, 구조 설계 (Human 80%)", "Step 2: Direct (지시)", "- 맥락(Context) 제공 및 제약 조건 설정 (Human 50% : AI 50%)", "Step 3: Audit (감리/승인)", "- 단순 디버깅/수정: AI Agent (Self-Healing)", "- 최종 품질 승인 및 책임(Accountability): Human 100%", "", "Message: AI는 실행하지만, 책임은 인간이 집니다."],
                  "일하는 방식도 바뀝니다. 단순 수정은 AI에게 맡기세요(Self-Healing). 하지만 최종적으로 이 코드를 서비스에 배포할지 결정하고, 그 결과에 책임지는 것은 오직 인간만이 할 수 있는 영역입니다.")

# 13. Conclusion
add_content_slide(prs, "결론: Be the Director", 
                  ["AI는 경쟁자가 아닙니다.", "- 당신이 고용한 가장 똑똑한 '비서'이자 '에이전트'입니다.", "", "\"Be the Director of your Code.\""],
                  "AI는 경쟁자가 아닙니다. 우리가 관리해야 할 유능한 팀원입니다. 이제 AI라는 엔진을 통해 여러분만의 소프트웨어를 '디자인' 하십시오.")

# 14. End
add_content_slide(prs, "Q & A", ["Thank You.", "", "질의 응답"], "경청해 주셔서 감사합니다.")

# 저장
file_path_updated = "Developer_Evolution_Agent_Ver.pptx"
prs.save(file_path_updated)
file_path_updated